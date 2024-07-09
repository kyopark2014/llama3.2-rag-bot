import json
import boto3
import os
import time
import datetime
import PyPDF2
import csv
import traceback
import re
import base64

from io import BytesIO
from urllib import parse
from botocore.config import Config

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferWindowMemory

from langchain_community.docstore.document import Document
from langchain_community.vectorstores.faiss import FAISS
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.embeddings import BedrockEmbeddings
from multiprocessing import Process, Pipe

from langchain_core.prompts import MessagesPlaceholder, ChatPromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_aws import ChatBedrock
from langchain_core.prompts import PromptTemplate

from langchain.agents import tool
from langchain.agents import AgentExecutor, create_react_agent
from bs4 import BeautifulSoup
from pytz import timezone
from langchain_community.tools.tavily_search import TavilySearchResults
from opensearchpy import OpenSearch

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
callLogTableName = os.environ.get('callLogTableName')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
enableReference = os.environ.get('enableReference', 'false')
debugMessageMode = os.environ.get('debugMessageMode', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
doc_prefix = s3_prefix+'/'
speech_prefix = 'speech/'
LLM_for_chat = json.loads(os.environ.get('LLM_for_chat'))
LLM_embedding = json.loads(os.environ.get('LLM_embedding'))
priorty_search_embedding = json.loads(os.environ.get('priorty_search_embedding'))
enalbeParentDocumentRetrival = os.environ.get('enalbeParentDocumentRetrival')
enableHybridSearch = os.environ.get('enableHybridSearch')

selected_chat = 0
selected_multimodal = 0
selected_embedding = 0
selected_ps_embedding = 0
separated_chat_history = os.environ.get('separated_chat_history')

useParallelRAG = os.environ.get('useParallelRAG', 'true')
roleArn = os.environ.get('roleArn')
top_k = int(os.environ.get('numberOfRelevantDocs', '8'))
MSG_LENGTH = 100
MSG_HISTORY_LENGTH = 20
speech_generation = True
history_length = 0
token_counter_history = 0

minDocSimilarity = 200
projectName = os.environ.get('projectName')

# api key to get weather information in agent
secretsmanager = boto3.client('secretsmanager')
try:
    get_weather_api_secret = secretsmanager.get_secret_value(
        SecretId=f"openweathermap-{projectName}"
    )
    #print('get_weather_api_secret: ', get_weather_api_secret)
    secret = json.loads(get_weather_api_secret['SecretString'])
    #print('secret: ', secret)
    weather_api_key = secret['weather_api_key']

except Exception as e:
    raise e
   
# api key to use LangSmith
langsmith_api_key = ""
try:
    get_langsmith_api_secret = secretsmanager.get_secret_value(
        SecretId=f"langsmithapikey-{projectName}"
    )
    #print('get_langsmith_api_secret: ', get_langsmith_api_secret)
    secret = json.loads(get_langsmith_api_secret['SecretString'])
    #print('secret: ', secret)
    langsmith_api_key = secret['langsmith_api_key']
    langchain_project = secret['langchain_project']
except Exception as e:
    raise e

if langsmith_api_key:
    os.environ["LANGCHAIN_API_KEY"] = langsmith_api_key
    os.environ["LANGCHAIN_TRACING_V2"] = "true"
    os.environ["LANGCHAIN_PROJECT"] = langchain_project
    
# api key to use Tavily Search
tavily_api_key = ""
try:
    get_tavily_api_secret = secretsmanager.get_secret_value(
        SecretId=f"tavilyapikey-{projectName}"
    )
    #print('get_tavily_api_secret: ', get_tavily_api_secret)
    secret = json.loads(get_tavily_api_secret['SecretString'])
    #print('secret: ', secret)
    tavily_api_key = secret['tavily_api_key']
except Exception as e: 
    raise e

if tavily_api_key:
    os.environ["TAVILY_API_KEY"] = tavily_api_key
    
# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

map_chain = dict() 

def get_chat():
    global selected_chat
    profile = LLM_for_chat[selected_chat]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_gen_len": 1024,  
        "top_p": 0.9, 
        "temperature": 0.1
    }    
    # print('parameters: ', parameters)

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters
    )    
    
    selected_chat = selected_chat + 1
    if selected_chat == len(LLM_for_chat):
        selected_chat = 0
    
    return chat

def get_embedding():
    global selected_embedding
    profile = LLM_embedding[selected_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'Embedding: {selected_embedding}, bedrock_region: {bedrock_region}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_embedding = selected_embedding + 1
    if selected_embedding == len(LLM_embedding):
        selected_embedding = 0
    
    return bedrock_embedding

def get_ps_embedding():
    global selected_ps_embedding
    profile = priorty_search_embedding[selected_ps_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_ps_embedding: {selected_ps_embedding}, bedrock_region: {bedrock_region}, model_id: {model_id}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_ps_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_ps_embedding = selected_ps_embedding + 1
    if selected_ps_embedding == len(priorty_search_embedding):
        selected_ps_embedding = 0
    
    return bedrock_ps_embedding

def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)

def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False

def general_conversation(connectionId, requestId, chat, query):
    global time_for_inference, history_length, token_counter_history    
    time_for_inference = history_length = token_counter_history = 0
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
    
    if isKorean(query)==True :
        system = (
            "다음의 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다. 답변은 한국어로 합니다."
        )
    else: 
        system = (
            "Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor."
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "history": history,
                "input": query,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
                            
        msg = stream.content
        print('msg: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
        
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
        
    return msg
    
def translate_text(chat, text):
    global time_for_inference
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
        
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    if isKorean(text)==False :
        input_language = "English"
        output_language = "Korean"
    else:
        input_language = "Korean"
        output_language = "English"
                        
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':          
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
    
    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def check_grammer(chat, text):
    global time_for_inference
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
        
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장의 오류를 찾아서 설명하고, 오류가 수정된 문장을 답변 마지막에 추가하여 주세요. 답변은 한국어로 합니다."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Find the error in the sentence and explain it, and add the corrected sentence at the end of your answer."
        )
        
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        msg = result.content
        print('result of grammer correction: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':          
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
    
    return msg

def get_summary(chat, docs):    
    text = ""
    for doc in docs:
        text = text + doc
    
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장을 요약해서 500자 이내의 한국어로 자세히 설명하세오."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Write a concise summary within 500 characters detaily."
        )
    
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        summary = result.content
        print('result of summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def revise_question(connectionId, requestId, chat, query):    
    global history_length, token_counter_history    
    history_length = token_counter_history = 0
        
    if isKorean(query)==True :      
        system = (
            ""
        )  
        human = """이전 대화를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다. 결과는 <result> tag를 붙여주세요. 답변은 한국어로 합니다.
        
        <question>            
        {question}
        </question>"""
        
    else: 
        system = (
            ""
        )
        human = """Rephrase the follow up <question> to be a standalone question. Put it in <result> tags.
        <question>            
        {question}
        </question>"""
            
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "history": history,
                "question": query,
            }
        )
        generated_question = result.content
        
        revised_question = generated_question[generated_question.find('<result>')+8:len(generated_question)-9] # remove <result> tag                   
        print('revised_question: ', revised_question)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")
            
    return revised_question    
    # return revised_question.replace("\n"," ")

def query_using_RAG_context(connectionId, requestId, chat, context, revised_question):    
    if isKorean(revised_question)==True:
        system = (
            """다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다. 답변은 한국어로 합니다.
            
            <context>
            {context}
            </context>"""
        )
    else: 
        system = (
            """Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
                   
    chain = prompt | chat
    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "context": context,
                "input": revised_question,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
        print('msg: ', msg)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    return msg
    
# load documents from s3 
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt' or file_type == 'md':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
            
        import docx
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load csv documents from s3
def load_csv_document(path, doc_prefix, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+doc_prefix+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs
    
def load_chat_history(userId, allowTime):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text' and text and msg:
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId):    
    msg_proceeding = {
        'request_id': requestId,
        'msg': 'Proceeding...',
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            # print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

def get_ps_embedding():
    global selected_ps_embedding
    profile = priorty_search_embedding[selected_ps_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_ps_embedding: {selected_ps_embedding}, bedrock_region: {bedrock_region}, model_id: {model_id}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_ps_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_ps_embedding = selected_ps_embedding + 1
    if selected_ps_embedding == len(priorty_search_embedding):
        selected_ps_embedding = 0
    
    return bedrock_ps_embedding

def priority_search(query, relevant_docs, minSimilarity):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        # print('doc: ', doc)
        if 'translated_excerpt' in doc['metadata'] and doc['metadata']['translated_excerpt']:
            content = doc['metadata']['translated_excerpt']
        else:
            content = doc['metadata']['excerpt']
            
        print('content: ', content)
        
        excerpts.append(
            Document(
                page_content=content,
                metadata={
                    'name': doc['metadata']['title'],
                    'order':i,
                }
            )
        )  
    # print('excerpts: ', excerpts)

    embeddings = get_ps_embedding()
    vectorstore_confidence = FAISS.from_documents(
        excerpts,  # documents
        embeddings  # embeddings
    )            
    rel_documents = vectorstore_confidence.similarity_search_with_score(
        query=query,
        # k=top_k
        k=len(relevant_docs)
    )

    docs = []
    for i, document in enumerate(rel_documents):
        # print(f'## Document(priority_search) {i+1}: {document}')

        order = document[0].metadata['order']
        name = document[0].metadata['name']
        assessed_score = document[1]
        print(f"{order} {name}: {assessed_score}")

        relevant_docs[order]['assessed_score'] = int(assessed_score)

        if assessed_score < minSimilarity:
            docs.append(relevant_docs[order])    
    # print('selected docs: ', docs)

    return docs

@tool 
def get_book_list(keyword: str) -> str:
    """
    Search book list by keyword and then return book list
    keyword: search keyword
    return: book list
    """
    
    keyword = keyword.replace('\'','')

    answer = ""
    url = f"https://search.kyobobook.co.kr/search?keyword={keyword}&gbCode=TOT&target=total"
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, "html.parser")
        prod_info = soup.find_all("a", attrs={"class": "prod_info"})
        
        if len(prod_info):
            answer = "추천 도서는 아래와 같습니다.\n"
            
        for prod in prod_info[:5]:
            title = prod.text.strip().replace("\n", "")       
            link = prod.get("href")
            answer = answer + f"{title}, URL: {link}\n\n"
    
    return answer
    
@tool
def get_current_time(format: str=f"%Y-%m-%d %H:%M:%S")->str:
    """Returns the current date and time in the specified format"""
    # f"%Y-%m-%d %H:%M:%S"
    
    format = format.replace('\'','')
    timestr = datetime.datetime.now(timezone('Asia/Seoul')).strftime(format)
    # print('timestr:', timestr)
    
    return timestr

@tool
def get_weather_info(city: str) -> str:
    """
    Search weather information by city name and then return weather statement.
    city: the english name of city to search
    return: weather statement
    """    
    
    city = city.replace('\n','')
    city = city.replace('\'','')
    
    chat = get_chat()
                
    if isKorean(city):
        place = traslation(chat, city, "Korean", "English")
        print('city (translated): ', place)
    else:
        place = city
        city = traslation(chat, city, "English", "Korean")
        print('city (translated): ', city)
        
    print('place: ', place)
    
    weather_str: str = f"{city}에 대한 날씨 정보가 없습니다."
    if weather_api_key: 
        apiKey = weather_api_key
        lang = 'en' 
        units = 'metric' 
        api = f"https://api.openweathermap.org/data/2.5/weather?q={place}&APPID={apiKey}&lang={lang}&units={units}"
        # print('api: ', api)
                
        try:
            result = requests.get(api)
            result = json.loads(result.text)
            print('result: ', result)
        
            if 'weather' in result:
                overall = result['weather'][0]['main']
                current_temp = result['main']['temp']
                min_temp = result['main']['temp_min']
                max_temp = result['main']['temp_max']
                humidity = result['main']['humidity']
                wind_speed = result['wind']['speed']
                cloud = result['clouds']['all']
                
                weather_str = f"{city}의 현재 날씨의 특징은 {overall}이며, 현재 온도는 {current_temp}도 이고, 최저온도는 {min_temp}도, 최고 온도는 {max_temp}도 입니다. 현재 습도는 {humidity}% 이고, 바람은 초당 {wind_speed} 미터 입니다. 구름은 {cloud}% 입니다."
                #weather_str = f"Today, the overall of {city} is {overall}, current temperature is {current_temp} degree, min temperature is {min_temp} degree, highest temperature is {max_temp} degree. huminity is {humidity}%, wind status is {wind_speed} meter per second. the amount of cloud is {cloud}%."            
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                    
            # raise Exception ("Not able to request to LLM")    
        
    print('weather_str: ', weather_str)                            
    return weather_str

@tool
def search_by_tavily(keyword: str) -> str:
    """
    Search general information by keyword and then return the result as a string.
    keyword: search keyword
    return: the information of keyword
    """    
    
    answer = ""    
    if tavily_api_key:
        keyword = keyword.replace('\'','')
        
        search = TavilySearchResults(k=3)
                    
        output = search.invoke(keyword)
        print('tavily output: ', output)
        
        for result in output:
            print('result: ', result)
            if result:
                content = result.get("content")
                url = result.get("url")
            
                answer = answer + f"{content}, URL: {url}\n"
        
    return answer

@tool    
def search_by_opensearch(keyword: str) -> str:
    """
    Search technical information by keyword and then return the result as a string.
    keyword: search keyword
    return: the technical information of keyword
    """    
    
    print('keyword: ', keyword)
    keyword = keyword.replace('\'','')
    keyword = keyword.replace('|','')
    keyword = keyword.replace('\n','')
    print('modified keyword: ', keyword)
    
    bedrock_embedding = get_embedding()
        
    vectorstore_opensearch = OpenSearchVectorSearch(
        index_name = "idx-*", # all
        is_aoss = False,
        ef_search = 1024, # 512(default)
        m=48,
        #engine="faiss",  # default: nmslib
        embedding_function = bedrock_embedding,
        opensearch_url=opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
    ) 
    
    answer = ""
    top_k = 2
    
    if enalbeParentDocumentRetrival == 'true': # parent/child chunking
        result = vectorstore_opensearch.similarity_search_with_score(
            query = keyword,
            k = top_k*2,  # use double
            pre_filter={"doc_level": {"$eq": "child"}}
        )
        print('result: ', result)
                
        relevant_documents = []
        docList = []
        for re in result:
            if 'parent_doc_id' in re[0].metadata:
                parent_doc_id = re[0].metadata['parent_doc_id']
                doc_level = re[0].metadata['doc_level']
                print(f"doc_level: {doc_level}, parent_doc_id: {parent_doc_id}")
                        
                if doc_level == 'child':
                    if parent_doc_id in docList:
                        print('duplicated!')
                    else:
                        relevant_documents.append(re)
                        docList.append(parent_doc_id)
                        
                        if len(relevant_documents)>=top_k:
                            break
                        
        for i, document in enumerate(relevant_documents):
            #print(f'## Document(opensearch-vector) {i+1}: {document}')
            
            parent_doc_id = document[0].metadata['parent_doc_id']
            doc_level = document[0].metadata['doc_level']
            print(f"child: parent_doc_id: {parent_doc_id}, doc_level: {doc_level}")
            
            excerpt, uri = get_parent_content(parent_doc_id)
            
            print(f"parent_doc_id: {parent_doc_id}, doc_level: {doc_level}, uri: {uri}, content: {excerpt}")
            
            answer = answer + f"{excerpt}, URL: {uri}\n\n"
    else: 
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = keyword,
            k = top_k,
        )

        for i, document in enumerate(relevant_documents):
            #print(f'## Document(opensearch-vector) {i+1}: {document}')
            
            excerpt = document[0].page_content        
            uri = document[0].metadata['uri']
                            
            answer = answer + f"{excerpt}, URL: {uri}\n\n"
    
    return answer

# define tools
tools = [get_current_time, get_book_list, get_weather_info, search_by_tavily, search_by_opensearch]        

def get_react_prompt_template(): # (hwchase17/react) https://smith.langchain.com/hub/hwchase17/react
    # Get the react prompt template    
    return PromptTemplate.from_template("""다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.

사용할 수 있는 tools은 아래와 같습니다:

{tools}

다음의 format을 사용하세요.:

Question: 답변하여야 할 input question 
Thought: you should always think about what to do. 
Action: 해야 할 action로서 [{tool_names}]에서 tool의 name만을 가져옵니다. 
Action Input: action의 input
Observation: action의 result
... (Thought/Action/Action Input/Observation을 5번 반복 할 수 있습니다.)
Thought: 나는 이제 Final Answer를 알고 있습니다. 
Final Answer: original input에 대한 Final Answer

너는 Human에게 해줄 응답이 있거나, Tool을 사용하지 않아도 되는 경우에, 다음 format을 사용하세요.:
'''
Thought: Tool을 사용해야 하나요? No
Final Answer: [your response here]
'''

Begin!

Question: {input}
Thought:{agent_scratchpad}
""")
        
def run_agent_react(connectionId, requestId, chat, query):
    prompt_template = get_react_prompt_template()
    print('prompt_template: ', prompt_template)
    
    #from langchain import hub
    #prompt_template = hub.pull("hwchase17/react")
    #print('prompt_template: ', prompt_template)
    
     # create agent
    isTyping(connectionId, requestId)
    agent = create_react_agent(chat, tools, prompt_template)
    
    agent_executor = AgentExecutor(
        agent=agent, 
        tools=tools, 
        verbose=True, 
        handle_parsing_errors=True,
        max_iterations = 5
    )
    
    # run agent
    response = agent_executor.invoke({"input": query})
    print('response: ', response)

    # streaming    
    msg = readStreamMsg(connectionId, requestId, response['output'])

    msg = response['output']
    print('msg: ', msg)
            
    return msg

def run_agent_react_chat_using_revised_question(connectionId, requestId, chat, query):
    # revise question
    revised_question = revise_question(connectionId, requestId, chat, query)     
    print('revised_question: ', revised_question)  
        
    # get template based on react 
    prompt_template = get_react_prompt_template()
    print('prompt_template: ', prompt_template)
    
    # create agent
    isTyping(connectionId, requestId)
    agent = create_react_agent(chat, tools, prompt_template)
    
    agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True, handle_parsing_errors=True)
    
    # run agent
    response = agent_executor.invoke({"input": revised_question})
    print('response: ', response)
    
    # streaming
    msg = readStreamMsg(connectionId, requestId, response['output'])

    msg = response['output']
    print('msg: ', msg)
            
    return msg

def get_react_chat_prompt_template():
    # Get the react prompt template

    return PromptTemplate.from_template("""다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.

사용할 수 있는 tools은 아래와 같습니다:

{tools}

다음의 format을 사용하세요.:

Question: 답변하여야 할 input question 
Thought: you should always think about what to do. 
Action: 해야 할 action로서 [{tool_names}]에서 tool의 name만을 가져옵니다. 
Action Input: action의 input
Observation: action의 result
... (Thought/Action/Action Input/Observation을 5번 반복 할 수 있습니다.)
Thought: 나는 이제 Final Answer를 알고 있습니다. 
Final Answer: original input에 대한 Final Answer

너는 Human에게 해줄 응답이 있거나, Tool을 사용하지 않아도 되는 경우에, 다음 format을 사용하세요.:
'''
Thought: Tool을 사용해야 하나요? No
Final Answer: [your response here]
'''

Begin!

Previous conversation history:
{chat_history}

New input: {input}
Thought:{agent_scratchpad}
""")
    
def run_agent_react_chat(connectionId, requestId, chat, query):
    # get template based on react 
    prompt_template = get_react_chat_prompt_template()
    print('prompt_template: ', prompt_template)
    
    # create agent
    isTyping(connectionId, requestId)
    agent = create_react_agent(chat, tools, prompt_template)
    
    agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True, handle_parsing_errors=True)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
    
    # run agent
    response = agent_executor.invoke({
        "input": query,
        "chat_history": history
    })
    print('response: ', response)
    
    # streaming
    msg = readStreamMsg(connectionId, requestId, response['output'])

    msg = response['output']
    print('msg: ', msg)
            
    return msg

def get_reference(docs):
    reference = "\n\nFrom\n"
    for i, doc in enumerate(docs):
        if doc['metadata']['translated_excerpt']:
            excerpt = str(doc['metadata']['excerpt']+'  [번역]'+doc['metadata']['translated_excerpt']).replace('"',"") 
        else:
            excerpt = str(doc['metadata']['excerpt']).replace('"'," ")
            
        excerpt = excerpt.replace('\n','\\n')            
                
        if doc['rag_type'][:10] == 'opensearch':
            #print(f'## Document(get_reference) {i+1}: {doc}')
                
            page = ""
            if "document_attributes" in doc['metadata']:
                if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                    page = doc['metadata']['document_attributes']['_excerpt_page_number']
            uri = doc['metadata']['source']
            name = doc['metadata']['title']
            #print('opensearch page: ', page)

            if page:                
                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
            else:
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                                               
    return reference

def get_documents_from_opensearch(vectorstore_opensearch, query, top_k):
    result = vectorstore_opensearch.similarity_search_with_score(
        query = query,
        k = top_k*2,  
        pre_filter={"doc_level": {"$eq": "child"}}
    )
    print('result: ', result)
            
    relevant_documents = []
    docList = []
    for re in result:
        if 'parent_doc_id' in re[0].metadata:
            parent_doc_id = re[0].metadata['parent_doc_id']
            doc_level = re[0].metadata['doc_level']
            print(f"doc_level: {doc_level}, parent_doc_id: {parent_doc_id}")
                    
            if doc_level == 'child':
                if parent_doc_id in docList:
                    print('duplicated!')
                else:
                    relevant_documents.append(re)
                    docList.append(parent_doc_id)
                    
                    if len(relevant_documents)>=top_k:
                        break
                                
    # print('lexical query result: ', json.dumps(response))
    print('relevant_documents: ', relevant_documents)
    
    return relevant_documents

os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress = True,
    http_auth=(opensearch_account, opensearch_passwd),
    use_ssl = True,
    verify_certs = True,
    ssl_assert_hostname = False,
    ssl_show_warn = False,
)

def retrieve_docs_from_vectorstore(vectorstore_opensearch, query, top_k):
    print(f"query: {query}")

    relevant_docs = []
            
    # vector search (semantic) 
    if enalbeParentDocumentRetrival=='true':
        relevant_documents = get_documents_from_opensearch(vectorstore_opensearch, query, top_k)
                
    else:
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        #print('(opensearch score) relevant_documents: ', relevant_documents)

    for i, document in enumerate(relevant_documents):
        #print('document.page_content:', document.page_content)
        #print('document.metadata:', document.metadata)
        #print(f'## Document(opensearch-vector) {i+1}: {document}')

        name = document[0].metadata['name']
        # print('metadata: ', document[0].metadata)

        page = ""
        if "page" in document[0].metadata:
            page = document[0].metadata['page']
        uri = ""
        if "uri" in document[0].metadata:
            uri = document[0].metadata['uri']

        excerpt = document[0].page_content
        confidence = str(document[1])
        assessed_score = str(document[1])
        
        parent_doc_id = doc_level = ""            
        if enalbeParentDocumentRetrival == 'true':
            parent_doc_id = document[0].metadata['parent_doc_id']
            doc_level = document[0].metadata['doc_level']
            excerpt, name, uri, doc_level = get_parent_document(parent_doc_id) # use pareant document

        if page:
            print('page: ', page)
            doc_info = {
                "rag_type": 'opensearch-vector',
                "confidence": confidence,
                "metadata": {
                    "source": uri,
                    "title": name,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "document_attributes": {
                        "_excerpt_page_number": page
                    },
                    "parent_doc_id": parent_doc_id,
                    "doc_level": doc_level          
                },
                "assessed_score": assessed_score,
            }
        else:
            doc_info = {
                "rag_type": 'opensearch-vector',
                "confidence": confidence,
                "metadata": {
                    "source": uri,
                    "title": name,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "parent_doc_id": parent_doc_id,
                    "doc_level": doc_level 
                },
                "assessed_score": assessed_score,
            }
        relevant_docs.append(doc_info)
        
    return relevant_docs

def get_parent_content(parent_doc_id):
    if parent_doc_id:
        response = os_client.get(
            index="idx-rag", 
            id = parent_doc_id
        )
            
        source = response['_source']
        print('excerpt: ', source['text'])   
            
        metadata = source['metadata']    
        #print('name: ', metadata['name'])   
        print('uri: ', metadata['uri'])   
        #print('doc_level: ', metadata['doc_level']) 
                    
    return source['text'], metadata['uri']

def get_parent_document(doc):
    # print('doc: ', doc)
    if 'parent_doc_id' in doc['metadata']:
        parent_doc_id = doc['metadata']['parent_doc_id']
    
        if parent_doc_id:
            response = os_client.get(
                index="idx-rag", 
                id = parent_doc_id
            )
            
            #source = response['_source']
            # print('parent_doc: ', source['text'])   
            
            #metadata = source['metadata']    
            #print('name: ', metadata['name'])   
            #print('uri: ', metadata['uri'])   
            #print('doc_level: ', metadata['doc_level']) 
            
            print('text(before)', doc['metadata']['excerpt'])
            doc['metadata']['excerpt'] = response['_source']['text']
            print('text(after)', doc['metadata']['excerpt'])
        
    return doc

def lexical_search(query, top_k):
    relevant_docs = []
    
    # lexical search (keyword)
    min_match = 0
    if enableHybridSearch == 'true':
        query = {
            "query": {
                "bool": {
                    "must": [
                        {
                            "match": {
                                "text": {
                                    "query": query,
                                    "minimum_should_match": f'{min_match}%',
                                    "operator":  "or",
                                }
                            }
                        },
                    ],
                    "filter": [
                    ]
                }
            }
        }

        response = os_client.search(
            body=query,
            index="idx-*", # all
        )
        # print('lexical query result: ', json.dumps(response))
        
        docList = []   
        for i, document in enumerate(response['hits']['hits']):
            if i>=top_k: 
                break
                    
            excerpt = document['_source']['text']
            #print(f'## Document(opensearch-keyword) {i+1}: {excerpt}')

            name = document['_source']['metadata']['name']
            # print('name: ', name)

            page = ""
            if "page" in document['_source']['metadata']:
                page = document['_source']['metadata']['page']
                    
            uri = ""
            if "uri" in document['_source']['metadata']:
                uri = document['_source']['metadata']['uri']
            # print('uri: ', uri)

            confidence = str(document['_score'])
            assessed_score = ""
            
            parent_doc_id = doc_level = ""
            if enalbeParentDocumentRetrival == 'true':
                if 'parent_doc_id' in document['_source']['metadata']:
                    parent_doc_id = document['_source']['metadata']['parent_doc_id']
                if 'doc_level' in document['_source']['metadata']:
                    doc_level = document['_source']['metadata']['doc_level']
                print(f"doc_level: {doc_level}, parent_doc_id: {parent_doc_id}")
                
            if page:
                print('page: ', page)
                doc_info = {
                    "rag_type": 'opensearch-keyword',
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": "",
                        "document_attributes": {
                            "_excerpt_page_number": page
                        },
                        "parent_doc_id": parent_doc_id,
                        "doc_level": doc_level  
                    },
                    "assessed_score": assessed_score,
                }
            else: 
                doc_info = {
                    "rag_type": 'opensearch-keyword',
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": "",
                        "parent_doc_id": parent_doc_id,
                        "doc_level": doc_level  
                    },
                    "assessed_score": assessed_score,
                }
            
            if parent_doc_id:  # parent doc
                if parent_doc_id in docList:  # check duplication partially                    
                    print('duplicated!')
                else:
                    relevant_docs.append(doc_info)
                    docList.append(parent_doc_id)
            else:  # child doc
                relevant_docs.append(doc_info)

    return relevant_docs

def vector_search(bedrock_embedding, query, top_k):
    print(f"query: {query}")
    
    vectorstore_opensearch = OpenSearchVectorSearch(
        index_name = "idx-*", # all
        is_aoss = False,
        ef_search = 1024, # 512(default)
        m=48,
        #engine="faiss",  # default: nmslib
        embedding_function = bedrock_embedding,
        opensearch_url=opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
    ) 

    relevant_docs = []
            
    # vector search (semantic) 
    if enalbeParentDocumentRetrival=='true':  # parent/child chunking
        result = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k*2,  # use double
            pre_filter={"doc_level": {"$eq": "child"}}
        )
        print('result of opensearch: ', result)
                
        relevant_documents = []
        docList = []
        for re in result:
            if 'parent_doc_id' in re[0].metadata:
                parent_doc_id = re[0].metadata['parent_doc_id']
                doc_level = re[0].metadata['doc_level']
                print(f"doc_level: {doc_level}, parent_doc_id: {parent_doc_id}")
                        
                if doc_level == 'child':
                    if parent_doc_id in docList:
                        print('duplicated!')
                    else:
                        relevant_documents.append(re)
                        docList.append(parent_doc_id)
                        
                        #if len(relevant_documents)>=top_k:
                        #    break
                
    else:  # single chunking
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        #print('(opensearch score) relevant_documents: ', relevant_documents)
        
    for i, document in enumerate(relevant_documents):
        #print('document.page_content:', document.page_content)
        #print('document.metadata:', document.metadata)
        #print(f'## Document(opensearch-vector) {i+1}: {document}')

        name = document[0].metadata['name']
        # print('metadata: ', document[0].metadata)

        page = ""
        if "page" in document[0].metadata:
            page = document[0].metadata['page']
        uri = ""
        if "uri" in document[0].metadata:
            uri = document[0].metadata['uri']

        excerpt = document[0].page_content
        confidence = str(document[1])
        assessed_score = str(document[1])
        
        parent_doc_id = doc_level = ""            
        if enalbeParentDocumentRetrival == 'true':
            parent_doc_id = document[0].metadata['parent_doc_id']
            doc_level = document[0].metadata['doc_level']

        if page:
            print('page: ', page)
            doc_info = {
                "rag_type": 'opensearch-vector',
                "confidence": confidence,
                "metadata": {
                    "source": uri,
                    "title": name,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "document_attributes": {
                        "_excerpt_page_number": page
                    },
                    "parent_doc_id": parent_doc_id,
                    "doc_level": doc_level  
                },
                "assessed_score": assessed_score,
            }
        else:
            doc_info = {
                "rag_type": 'opensearch-vector',
                "confidence": confidence,
                "metadata": {
                    "source": uri,
                    "title": name,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "parent_doc_id": parent_doc_id,
                    "doc_level": doc_level  
                },
                "assessed_score": assessed_score,
            }
        relevant_docs.append(doc_info)
            
    return relevant_docs

def get_answer_using_RAG(chat, text, search_type, connectionId, requestId, bedrock_embedding):
    global time_for_revise, time_for_rag, time_for_inference, time_for_priority_search, number_of_relevant_docs 
    time_for_revise = time_for_rag = time_for_inference = time_for_priority_search = number_of_relevant_docs = 0
    
    start_time_for_revise = time.time()
    
     # revise question
    revised_question = revise_question(connectionId, requestId, chat, text)     
    print('revised_question: ', revised_question)
    
    end_time_for_revise = time.time()
    time_for_revise = end_time_for_revise - start_time_for_revise
    print('processing time for revised question: ', time_for_revise)
    
    # retrieve relevant documents from RAG
    selected_relevant_docs = retrieve_docs_from_RAG(revised_question, connectionId, requestId, bedrock_embedding, search_type)
    
    # get context
    relevant_context = ""
    number_of_relevant_docs = len(selected_relevant_docs)    
    for document in selected_relevant_docs:
        content = document['metadata']['excerpt']
                
        relevant_context = relevant_context + content + "\n\n"
    # print('relevant_context: ', relevant_context)
    
    end_time_for_rag = time.time()
    time_for_rag = end_time_for_rag - end_time_for_revise
    print('processing time for RAG: ', time_for_rag)
        
    # query using RAG context
    msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, revised_question)

    reference = ""
    if len(selected_relevant_docs)>=1 and enableReference=='true':
        reference = get_reference(selected_relevant_docs)  

    end_time_for_inference = time.time()
    time_for_inference = end_time_for_inference - end_time_for_rag
    print('processing time for inference: ', time_for_inference)
    
    global relevant_length, token_counter_relevant_docs    
    if debugMessageMode=='true':   # extract chat history for debug
        relevant_length = len(relevant_context)
        token_counter_relevant_docs = chat.get_num_tokens(relevant_context)

    return msg, reference
    
def retrieve_docs_from_RAG(revised_question, connectionId, requestId, bedrock_embedding, search_type):
    # vector search
    rel_docs_vector_search = vector_search(bedrock_embedding=bedrock_embedding, query=revised_question, top_k=top_k)
    print(f'rel_docs (vector): '+json.dumps(rel_docs_vector_search))
    
    if search_type == 'hybrid':
        # lexical search
        rel_docs_lexical_search = lexical_search(revised_question, top_k)    
        print(f'rel_docs (lexical): '+json.dumps(rel_docs_lexical_search))
        relevant_docs = rel_docs_vector_search + rel_docs_lexical_search
    else:  # vector only
        relevant_docs = rel_docs_vector_search    
    
    # priority search
    global time_for_priority_search
    time_for_priority_search = 0    
    start_time_for_priority_search = time.time()
    
    selected_relevant_docs = []
    if len(relevant_docs)>=1:
        print('start priority search')
        selected_relevant_docs = priority_search(revised_question, relevant_docs, minDocSimilarity)
        print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))
        
    # update doc using parent
    contentList = []
    update_docs = []
    for doc in selected_relevant_docs:        
        doc = get_parent_document(doc) # use pareant document
        
        # print('excerpt: ', doc['metadata']['excerpt'])
        if doc['metadata']['excerpt'] in contentList:
            print('duplicated!')
            continue
        contentList.append(doc['metadata']['excerpt'])
        update_docs.append(doc)
        
        if len(update_docs)>=top_k:
            break
    
    print('update_docs:', json.dumps(update_docs))
    #for i, doc in enumerate(update_docs):
    #    print(f"#### relevant_docs ({i}): {json.dumps(update_docs)}")
                    
    end_time_for_priority_search = time.time() 
    time_for_priority_search = end_time_for_priority_search - start_time_for_priority_search
    print('processing time for priority search: ', time_for_priority_search)
    
    return update_docs
    
def traslation(chat, text, input_language, output_language):
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        # print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)
    
    rag_type = ""
    if 'rag_type' in jsonBody:
        if jsonBody['rag_type']:
            rag_type = jsonBody['rag_type']  # RAG type
            print('rag_type: ', rag_type)
    
    global enableReference
    global map_chain, memory_chain, debugMessageMode
                 
    # Multi-LLM
    profile = LLM_for_chat[selected_chat]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_chat: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
      
    chat = get_chat()    
    bedrock_embedding = get_embedding()

    # allocate memory
    if userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        allowTime = getAllowTime()
        load_chat_history(userId, allowTime)
        
    start = int(time.time())    

    msg = ""
    reference = ""
    isControlMsg = False
    token_counter_input = 0
    time_for_revise = time_for_rag = time_for_priority_search = time_for_inference = 0
    
    if type == 'text' and body[:11] == 'list models':
        isControlMsg = True
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)          

        sendResultMessage(connectionId, requestId, msg)  
    else:           
        text = body
        print('query: ', text)  
        querySize = len(text)
        textCount = len(text.split())
        print(f"query size: {querySize}, words: {textCount}")
        
        if type == 'text':
            if text == 'enableReference':
                enableReference = 'true'
                isControlMsg = True
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                isControlMsg = True
                msg  = "Reference is disabled"
            elif text == 'enableDebug':
                isControlMsg = True
                debugMessageMode = 'true'
                msg  = "Debug messages will be delivered to the client."
            elif text == 'disableDebug':
                isControlMsg = True
                debugMessageMode = 'false'
                msg = "Debug messages will not be delivered to the client."

            elif text == 'clearMemory':
                isControlMsg = True
                memory_chain.clear()
                map_chain[userId] = memory_chain
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:       
                if conv_type == 'normal':      # normal
                    msg = general_conversation(connectionId, requestId, chat, text)      
                
                elif conv_type == 'agent-react':
                    msg = run_agent_react(connectionId, requestId, chat, text)                
                elif conv_type == 'agent-react-chat':         
                    if separated_chat_history=='true': 
                        msg = run_agent_react_chat_using_revised_question(connectionId, requestId, chat, text)
                    else:
                        msg = run_agent_react_chat(connectionId, requestId, chat, text)
                            
                elif conv_type == 'qa-opensearch-vector':   # RAG - Vector
                    print(f'rag_type: {rag_type}')
                    search_type ='vector'
                    msg, reference = get_answer_using_RAG(chat, text, search_type, connectionId, requestId, bedrock_embedding)
                
                elif conv_type == 'qa-opensearch-hybrid':   # RAG - Hybrid
                    print(f'rag_type: {rag_type}')
                    search_type = 'hybrid'
                    msg, reference = get_answer_using_RAG(chat, text, search_type, connectionId, requestId, bedrock_embedding)
                        
                elif conv_type == 'translation':                    
                    msg = translate_text(chat, text)
                
                elif conv_type == 'grammar':                    
                    msg = check_grammer(chat, text)
                    
                # token counter
                if debugMessageMode=='true':
                    token_counter_input = chat.get_num_tokens(text)
                    token_counter_output = chat.get_num_tokens(msg)
                    print(f"token_counter: question: {token_counter_input}, answer: {token_counter_output}")
                    
                memory_chain.chat_memory.add_user_message(text)
                memory_chain.chat_memory.add_ai_message(msg)
                        
        elif type == 'document':
            isTyping(connectionId, requestId)
            
            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)
            
            if file_type == 'csv':
                docs = load_csv_document(path, doc_prefix, object)
                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
                        
            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+doc_prefix+parse.quote(object)
                            }
                        )
                    )
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                contexts = []
                for doc in docs[:2]:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
                
            else:
                msg = "uploaded file: "+object
                                                        
        sendResultMessage(connectionId, requestId, msg+reference)
        # print('msg+reference: ', msg+reference)

        elapsed_time = time.time() - start
        print("total run time(sec): ", elapsed_time)
                               
        item = {    # save dialog
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if debugMessageMode=='true' and isControlMsg==False: 
        statusMsg = f"\n\n[통계]\nRegion: {bedrock_region}\nModelId: {modelId}\n"
        if token_counter_input:
            statusMsg = statusMsg + f"Question: {str(len(text))}자 / {token_counter_input}토큰\nAnswer: {str(len(msg))}자 / {token_counter_output}토큰\n"
            
        if history_length:
            statusMsg = statusMsg + f"History: {str(history_length)}자 / {token_counter_history}토큰\n"
            
        statusMsg = statusMsg + f"Time(초): "
        if time_for_revise != 0:
            statusMsg = statusMsg + f"{time_for_revise:.2f}(Revise), "
        if time_for_rag != 0:
            statusMsg = statusMsg + f"{time_for_rag:.2f}(RAG), "
        if time_for_priority_search != 0:
            statusMsg = statusMsg + f"{time_for_priority_search:.2f}(Priority) "
        if time_for_inference != 0:
            statusMsg = statusMsg + f"{time_for_inference:.2f}(Inference), "
        statusMsg = statusMsg + f"{elapsed_time:.2f}(전체)"
            
        sendResultMessage(connectionId, requestId, msg+reference+statusMsg)

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                                        
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")

    return {
        'statusCode': 200
    }
